from dotenv import load_dotenv
load_dotenv()

from flask import Flask, render_template, request, redirect, url_for, flash, after_this_request
import google.generativeai as genai
import tempfile
from flask import send_file
from werkzeug.utils import secure_filename
import PyPDF2
import docx
import pptx
import re
import os

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "your_secret_key")  # Flash mesajları için bir secret key gerekli.

def save_api_key_to_env(api_key):
    env_path = os.path.join(os.path.dirname(__file__), '.env')
    lines = []
    if os.path.exists(env_path):
        with open(env_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
    key_exists = False
    for i, line in enumerate(lines):
        if line.startswith('GEMINI_API_KEY='):
            lines[i] = f'GEMINI_API_KEY={api_key}\n'
            key_exists = True
            break
    if not key_exists:
        lines.append(f'GEMINI_API_KEY={api_key}\n')
    with open(env_path, 'w', encoding='utf-8') as f:
        f.writelines(lines)

def get_gemini_api_key():
    return os.environ.get("GEMINI_API_KEY")

GEMINI_API_KEY = get_gemini_api_key()

if not GEMINI_API_KEY:
    # We will handle redirect in route instead of raising error here
    pass
else:
    genai.configure(api_key=GEMINI_API_KEY)

# Kullanılabilir modelleri listeleyin ve desteklenen bir modeli seçin
models = genai.list_models()
#models = list(models)  # Generator nesnesini listeye dönüştürün
#print(f"Kullanılabilir Modeller: {models}")
if not models:
    raise ValueError("Hiçbir desteklenen model bulunamadı. Lütfen API anahtarınızı kontrol edin.")
# Varsayılan olarak gemini-1.5-flash modelini seçiyoruz
selected_model_name = "models/gemini-2.0-flash"
model = genai.GenerativeModel(selected_model_name)

# Global dictionary to store chat sessions
chat_sessions = {}

def convert_bold(text):
    # Replace **text** with <strong>text</strong>
    return re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)


def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_path, filename):
    ext = filename.rsplit('.', 1)[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext in ["doc", "docx"]:
        return extract_text_from_docx(file_path)
    elif ext in ["ppt", "pptx"]:
        return extract_text_from_pptx(file_path)
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    else:
        return ""

import os
import traceback

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return "No file part", 400
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400
    filename = secure_filename(file.filename)

    # Get prompt selections from form data, default to False if not present
    prompt1_selected = request.form.get('prompt1', 'false').lower() == 'true'
    prompt2_selected = request.form.get('prompt2', 'false').lower() == 'true'

    if not prompt1_selected and not prompt2_selected:
        return "No prompt selected. Please select at least one prompt.", 400

    try:
        tmp = tempfile.NamedTemporaryFile(delete=False)
        try:
            file.save(tmp.name)
            text = extract_text(tmp.name, filename)
        finally:
            tmp.close()
        os.unlink(tmp.name)  # Delete temp file manually

        if not text.strip():
            return "Unsupported file type or empty file", 400

        combined_result = ""

        if prompt1_selected:
            prompt1 = (
                "Bu metnin tamamını her şeye değindiğinden emin olarak anlam karmaşasını ve yanlış yazımları düzelterek "
                "konu anlatan bir hocanın anlatım tarzıyla baştan yazar mısın?\n\n"
                f"{text}"
            )
            response1 = model.generate_content(prompt1)
            rewritten_text = response1.text.strip()
            combined_result += "=== Geliştirilmiş Anlatım ===\n" + rewritten_text + "\n\n"

        if prompt2_selected:
            prompt2 = (
                "bu dosyadan 20 tane zor seviyede (TUS tarzında) soru hazırlar mısın internet araması yapabilirsin her soru 5 seçenekten oluşsun önce tüm soruları sor sonra tüm cevapları ver sonra tüm çözümleri yazar mısın?\n\n"
                f"{text}"
            )
            response2 = model.generate_content(prompt2)
            soru_text = response2.text.strip()
            combined_result += "=== Dersten Sorular ===\n" + soru_text + "\n\n"

        import docx
        doc = docx.Document()

        def add_bold_paragraph(doc, text):
            # This function adds a paragraph to the doc with **text** parts bolded
            pattern = re.compile(r'\*\*(.+?)\*\*')
            paragraph = doc.add_paragraph()
            last_end = 0
            for match in pattern.finditer(text):
                # Add text before the match
                if match.start() > last_end:
                    paragraph.add_run(text[last_end:match.start()])
                # Add bold text
                bold_run = paragraph.add_run(match.group(1))
                bold_run.bold = True
                last_end = match.end()
            # Add remaining text after last match
            if last_end < len(text):
                paragraph.add_run(text[last_end:])
            return paragraph

        for line in combined_result.splitlines():
            add_bold_paragraph(doc, line)

        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as out_file:
            doc.save(out_file.name)
            out_file_path = out_file.name

        # Change download_name to uploaded filename with .docx extension
        base_filename = os.path.splitext(filename)[0]
        download_filename = base_filename + ".docx"
        response = send_file(out_file_path, as_attachment=True, download_name=download_filename, mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

        @after_this_request
        def remove_file(response):
            try:
                os.unlink(out_file_path)
            except Exception as e:
                app.logger.error(f"Error deleting temp file: {e}")
            return response

        return response
    except Exception as e:
        error_message = f"An error occurred: {str(e)}"
        app.logger.error(f"{error_message}\n{traceback.format_exc()}")
        return error_message, 500
    
@app.route('/test', methods=['POST'])
def test():
    if 'file' not in request.files:
        return "No file part", 400
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400
    filename = secure_filename(file.filename)
    try:
        tmp = tempfile.NamedTemporaryFile(delete=False)
        try:
            file.save(tmp.name)
            text = extract_text(tmp.name, filename)
        finally:
            tmp.close()
        os.unlink(tmp.name)  # Delete temp file manually

        if not text.strip():
            return "Unsupported file type or empty file", 400

        prompt2 = ("bu dosyadan 20 tane zor seviyede soru hazırlar mısın internet araması yapabilirsin her soru 5 seçenekten oluşsun önce tüm soruları sor sonra tüm cevapları ver sonra tüm çözümleri yazar mısın?\n\n"f"{text}")

        response = model.generate_content(prompt2)
        soru_text = response.text.strip()

        import docx
        doc = docx.Document()
        for line in soru_text.splitlines():
            doc.add_paragraph(line)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as out_file:
            doc.save(out_file.name)
            out_file_path = out_file.name

        # Change download_name to uploaded filename with .docx extension
        base_filename = os.path.splitext(filename)[0]
        download_filename = base_filename + ".docx"
        response = send_file(out_file_path, as_attachment=True, download_name=download_filename, mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

        @after_this_request
        def remove_file(response):
            try:
                os.unlink(out_file_path)
            except Exception as e:
                app.logger.error(f"Error deleting temp file: {e}")
            return response

        return response
    except Exception as e:
        error_message = f"An error occurred: {str(e)}"
        app.logger.error(f"{error_message}\n{traceback.format_exc()}")
        return error_message, 500

@app.route('/')
def index():
    gemini_key = get_gemini_api_key()
    if not gemini_key:
        return redirect(url_for('set_api_key'))
    return render_template('index.html')

@app.route('/set_api_key', methods=['GET', 'POST'])
def set_api_key():
    if request.method == 'POST':
        api_key = request.form.get('api_key')
        if not api_key:
            flash("API anahtarı boş olamaz.", "error")
            return redirect(url_for('set_api_key'))
        save_api_key_to_env(api_key)
        # Reload environment variables
        load_dotenv(override=True)
        global GEMINI_API_KEY
        GEMINI_API_KEY = get_gemini_api_key()
        genai.configure(api_key=GEMINI_API_KEY)
        flash("API anahtarı kaydedildi.", "success")
        return redirect(url_for('index'))
    return render_template('set_api_key.html')


if __name__ == '__main__':
    app.run(debug=True)
